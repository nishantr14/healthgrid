# Builds the 3-slide Demo Day deck (New Delhi, before MPs — exactly 3 slides,
# 3 minutes) in a clean light theme: white background, near-black ink, deep
# HealthGrid pink as the only accent. Fonts are Segoe UI throughout so the
# deck renders identically on any Windows machine — nothing to embed.
#   1. Problem & Solution        — headline stats + existing-vs-HealthGrid table
#   2. Key Features & Innovation — command-center screenshot + feature cards
#   3. Impact & Scalability      — measured numbers + scale rails + roadmap + ask
# Run: python scripts/demoday-deck-light.py [out.pptx]
import sys
from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import MSO_ANCHOR, PP_ALIGN
from pptx.util import Emu, Inches, Pt

OUT = sys.argv[1] if len(sys.argv) > 1 else r"docs/HealthGrid-Demoday-3Slides-Light.pptx"
SHOTS = r"C:\Users\nisha\AppData\Local\Temp\claude\C--bwa\930f78d5-b03d-4478-8177-d9aa79f4a0bb\scratchpad"

INK = RGBColor(0x1B, 0x16, 0x30)
GREY = RGBColor(0x5C, 0x56, 0x6E)
MUTED = RGBColor(0x8A, 0x84, 0x96)
ACCENT = RGBColor(0xD6, 0x33, 0x6C)
ACCENT_TINT = RGBColor(0xFC, 0xEB, 0xF2)
CARD = RGBColor(0xF6, 0xF4, 0xF8)
LINE = RGBColor(0xE6, 0xE2, 0xEC)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
FONT = "Segoe UI"

pres = Presentation()
pres.slide_width = Inches(13.333)
pres.slide_height = Inches(7.5)
BLANK = pres.slide_layouts[6]


def add_slide():
    return pres.slides.add_slide(BLANK)


def rect(slide, x, y, w, h, fill, line_color=None, line_w=0.75, round_=0.0):
    shape = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE if round_ else MSO_SHAPE.RECTANGLE,
        Inches(x), Inches(y), Inches(w), Inches(h),
    )
    if round_:
        shape.adjustments[0] = round_
    if fill is None:
        shape.fill.background()
    else:
        shape.fill.solid()
        shape.fill.fore_color.rgb = fill
    if line_color is None:
        shape.line.fill.background()
    else:
        shape.line.color.rgb = line_color
        shape.line.width = Pt(line_w)
    shape.shadow.inherit = False
    return shape


def text(slide, x, y, w, h, runs, *, size=11, color=INK, bold=False, align=PP_ALIGN.LEFT,
         anchor=MSO_ANCHOR.TOP, leading=1.0, space_before=0):
    """runs: string, or list of paragraphs, each a string or list of (text, overrides)."""
    tb = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    tf = tb.text_frame
    tf.word_wrap = True
    tf.vertical_anchor = anchor
    tf.margin_left = tf.margin_right = tf.margin_top = tf.margin_bottom = 0
    if isinstance(runs, str):
        runs = [runs]
    for i, para in enumerate(runs):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = align
        if leading != 1.0:
            p.line_spacing = leading
        if space_before and i > 0:
            p.space_before = Pt(space_before)
        if isinstance(para, str):
            para = [(para, {})]
        for run_text, over in para:
            r = p.add_run()
            r.text = run_text
            r.font.name = FONT
            r.font.size = Pt(over.get("size", size))
            r.font.bold = over.get("bold", bold)
            r.font.color.rgb = over.get("color", color)
    return tb


def pill(slide, label, x=0.7, y=0.48):
    w = 0.16 + len(label) * 0.082
    rect(slide, x, y, w, 0.34, ACCENT_TINT, round_=0.5)
    text(slide, x, y + 0.055, w, 0.24, label, size=10.5, color=ACCENT, bold=True, align=PP_ALIGN.CENTER)


def footer(slide, n):
    rect(slide, 0.7, 7.08, 11.93, 0.012, LINE)
    text(slide, 0.7, 7.16, 8.0, 0.25, "HealthGrid AI  ·  Team Noida Boys — Smart Health, Code for Communities", size=8.5, color=MUTED)
    text(slide, 11.63, 7.16, 1.0, 0.25, f"{n} / 3", size=8.5, color=MUTED, align=PP_ALIGN.RIGHT)


def accent_band(slide, y, h, runs, fill=ACCENT_TINT, color=INK, size=13):
    rect(slide, 0.7, y, 11.93, h, fill, round_=0.10)
    rect(slide, 0.7, y, 0.07, h, ACCENT)
    text(slide, 1.0, y, 11.4, h, runs, size=size, color=color, anchor=MSO_ANCHOR.MIDDLE, leading=1.08)


# ============================ Slide 1 ======================================
s = add_slide()
pill(s, "HEALTHGRID AI  ·  PROBLEM & SOLUTION")
text(s, 0.7, 1.0, 12.0, 0.65, "India's primary healthcare runs blind between monthly reports.", size=29, bold=True)
text(s, 0.7, 1.62, 12.0, 0.4,
     "Medicines, beds and staff for a district's health centres live in paper registers — the district finds out about a crisis about a month too late.",
     size=13, color=GREY)

stats = [
    ("25,000+", "PHCs on paper", "Medicines, beds and staff tracked in registers — no real-time visibility"),
    ("~30 days", "Reporting delay", "Before facility data reaches the district through monthly HMIS reporting"),
    ("0", "Early warnings", "Stock-outs and staffing gaps are found only after patients are turned away"),
]
sy = 2.42
for num, label, desc in stats:
    rect(s, 0.7, sy, 4.55, 1.16, CARD, LINE, round_=0.09)
    text(s, 0.98, sy + 0.12, 1.9, 0.55, num, size=26, color=ACCENT, bold=True)
    text(s, 2.95, sy + 0.13, 2.15, 0.3, label, size=11.5, bold=True)
    text(s, 2.95, sy + 0.44, 2.15, 0.68, desc, size=9, color=GREY, leading=1.05)
    sy += 1.31

text(s, 5.62, 2.42, 6.99, 0.3, "EXISTING SYSTEMS  VS  HEALTHGRID AI", size=10.5, color=ACCENT, bold=True)
rows = [
    ("Monthly reporting", "Live facility status"),
    ("Historical records", "Real-time risk score, 0–100"),
    ("Drug logistics records", "Early stock-out warnings"),
    ("Manual escalation", "AI-prioritized intervention queue"),
    ("Static dashboards", "Forecast → recommend → approve"),
]
ty, rh = 2.78, 0.66
rect(s, 5.62, ty, 6.99, rh * len(rows), WHITE, LINE, round_=0.0)
for i, (old, new) in enumerate(rows):
    y = ty + i * rh
    if i:
        rect(s, 5.62, y, 6.99, 0.012, LINE)
    text(s, 5.86, y, 3.1, rh, old, size=11.5, color=GREY, anchor=MSO_ANCHOR.MIDDLE)
    text(s, 9.06, y, 0.3, rh, "→", size=11.5, color=ACCENT, bold=True, anchor=MSO_ANCHOR.MIDDLE)
    text(s, 9.44, y, 3.0, rh, new, size=11.5, color=INK, bold=True, anchor=MSO_ANCHOR.MIDDLE)
rect(s, 5.62, ty, 0.05, rh * len(rows), ACCENT)

accent_band(s, 6.32, 0.62, [[
    ("The solution — HealthGrid AI:  ", {"bold": True, "color": ACCENT}),
    ("a real-time decision layer on the records districts already keep. HMIS records what happened. HealthGrid decides what to do next.", {}),
]])
footer(s, 1)

# ============================ Slide 2 ======================================
s = add_slide()
pill(s, "KEY FEATURES & INNOVATION")
text(s, 0.7, 1.0, 12.0, 0.65, "A real-time command center for every district decision.", size=29, bold=True)
text(s, 0.7, 1.62, 12.0, 0.4,
     "Deterministic engines score every facility — Gemini explains the numbers, listens to the frontline, and speaks its languages.",
     size=13, color=GREY)

pic = s.shapes.add_picture(SHOTS + r"\dash.png", Inches(0.7), Inches(2.2), height=Inches(4.62))
pic.line.color.rgb = LINE
pic.line.width = Pt(1)

cards = [
    ("Live Digital Twin", "15 facilities scored 0–100 on medicines, staffing, beds, surge and diagnostics — live on a district map."),
    ("Voice-First Frontline", "A health worker speaks in Hindi, Marathi or English — Gemini parses it and the district re-scores in seconds."),
    ("AI That Explains & Acts", "Root-cause analysis and guarded medicine transfers between facilities — approved in one click, fully audited."),
    ("WhatsApp Alerts + PDF Reports", "One-click alerts to frontline WhatsApp and a meeting-ready district report — a closed, acknowledged loop."),
]
cy = 2.2
for title, body in cards:
    rect(s, 6.65, cy, 5.98, 0.98, CARD, LINE, round_=0.10)
    rect(s, 6.65, cy, 0.06, 0.98, ACCENT)
    text(s, 6.93, cy + 0.11, 5.55, 0.28, title, size=12.5, bold=True)
    text(s, 6.93, cy + 0.42, 5.55, 0.5, body, size=10, color=GREY, leading=1.05)
    cy += 1.13
text(s, 6.93, cy + 0.02, 5.7, 0.5,
     "Also inside: weather & disaster stress mode  ·  deterministic engines, 88 unit tests  ·  Gemini + Firestore + Google Maps, on Cloud Run.",
     size=9.5, color=MUTED, leading=1.1)
footer(s, 2)

# ============================ Slide 3 ======================================
s = add_slide()
pill(s, "IMPACT & SCALABILITY")
text(s, 0.7, 1.0, 12.0, 0.65, "One district first. Built for 800+.", size=29, bold=True)
text(s, 0.7, 1.62, 12.0, 0.4,
     "A simulated month in Wardha district — every decision produced by the same engines and guardrails that run in the product.",
     size=13, color=GREY)

impact = [
    ("Stock-outs prevented", "Shortages are forecast days in advance and resolved before a single patient is turned away."),
    ("Guarded transfers", "Every redistribution is validated against safety guardrails — and approved by a human, never auto-applied."),
    ("Minimal new spend", "Demand is met by smarter redistribution of the stock the district already holds, not new purchases."),
    ("No manual data entry", "Runs on the records districts already keep, updated by voice from the frontline — no added paperwork."),
]
gx, gy, gw, gh, gap = 0.7, 2.42, 3.02, 1.72, 0.18
for i, (label, desc) in enumerate(impact):
    x = gx + (i % 2) * (gw + gap)
    y = gy + (i // 2) * (gh + gap)
    rect(s, x, y, gw, gh, CARD, LINE, round_=0.08)
    rect(s, x, y, gw, 0.06, ACCENT)
    text(s, x + 0.24, y + 0.22, gw - 0.48, 0.55, label, size=14, color=ACCENT, bold=True)
    text(s, x + 0.24, y + 0.62, gw - 0.48, 1.0, desc, size=10, color=GREY, leading=1.1)

rx = 7.0
text(s, rx, 2.42, 5.6, 0.28, "SCALE ON EXISTING RAILS", size=11, color=ACCENT, bold=True)
text(s, rx, 2.76, 5.62, 1.15, [
    "•  Wardha → Maharashtra → 800+ districts nationally — integrating with HMIS and DVDMS, the systems states already run",
    "•  Funded through NHM's Programme Implementation Plan — no new hardware, no new data entry, low thousands ₹/month per district",
], size=11, color=INK, leading=1.1, space_before=6)
text(s, rx, 4.18, 5.6, 0.28, "WHAT COMES NEXT", size=11, color=ACCENT, bold=True)
text(s, rx, 4.52, 5.62, 1.5, [
    "•  Offline-first frontline — voice and manual updates queue on the device and sync when the network returns",
    "•  Every regional language — each state rollout ships its own",
    "•  Readiness scores — operational & facility readiness metrics as new data streams come online",
], size=11, color=INK, leading=1.1, space_before=6)

accent_band(s, 6.32, 0.62, [[
    ("Our ask: sanction one pilot district.  ", {"bold": True, "color": WHITE}),
    ("Its own numbers will make the case to the state.", {"color": WHITE}),
]], fill=ACCENT, size=14)
footer(s, 3)

# Hanging indents for every bullet paragraph on slide 3.
for shape in s.shapes:
    if shape.has_text_frame:
        for p in shape.text_frame.paragraphs:
            if p.runs and p.runs[0].text.startswith("•"):
                pPr = p._p.get_or_add_pPr()
                pPr.set("marL", "171450")
                pPr.set("indent", "-171450")

pres.save(OUT)
print("WROTE", OUT, "slides:", len(pres.slides._sldIdLst))
