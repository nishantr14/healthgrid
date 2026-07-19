# Builds the 3-slide Demo Day deck (New Delhi, in front of MPs — exactly 3
# slides per the organizers' brief) from the team's winning v4_1 deck:
#   1. Problem & Solution   — the existing-systems vs HealthGrid comparison
#   2. Key Features & Innovation — command center + voice + WhatsApp cards
#   3. Impact & Scalability — measured impact numbers + scale-on-rails ask
# Run: python scripts/demoday-deck.py [out.pptx]
import sys
from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.util import Inches, Pt

SRC = r"C:\Users\nisha\Downloads\HealthGrid-Pitch-Final-v4_1.pptx"
OUT = sys.argv[1] if len(sys.argv) > 1 else r"docs/HealthGrid-Demoday-3Slides.pptx"

PINK = RGBColor(0xFF, 0x8A, 0xAF)
BODY = RGBColor(0xE0, 0xD6, 0xDE)
MUTED = RGBColor(0xB9, 0xAD, 0xC4)
PANEL_BG = RGBColor(0x14, 0x0B, 0x2A)
PANEL_LINE = RGBColor(0x7C, 0x6B, 0xAF)

pres = Presentation(SRC)


def by_text(slide, prefix):
    for sh in slide.shapes:
        if sh.has_text_frame and sh.text_frame.text.strip().startswith(prefix):
            return sh
    raise KeyError(prefix)


def set_text(shape, text):
    """Replace the shape's text keeping the first run's formatting."""
    p0 = shape.text_frame.paragraphs[0]
    p0.runs[0].text = text
    for r in p0.runs[1:]:
        r.text = ""
    for p in shape.text_frame.paragraphs[1:]:
        for r in p.runs:
            r.text = ""


# ================= Slide 1: Problem & Solution (base: THE INSIGHT) =========
s1 = pres.slides[2]
chip = by_text(s1, "THE INSIGHT")
set_text(chip, "HEALTHGRID AI · PROBLEM & SOLUTION")
chip.width = Inches(4.6)
for sh in s1.shapes:  # widen the chip pill behind the label
    if sh.name == "Shape 0":
        sh.width = Inches(3.85)
set_text(by_text(s1, "HMIS records"), "HMIS records what happened. HealthGrid AI decides what to do next.")
set_text(
    by_text(s1, "No rip-and-replace"),
    "The problem: 25,000+ PHCs on paper, a ~30-day reporting delay, zero early warnings — stock-outs surface only "
    "after patients are turned away. HealthGrid runs on the records districts already keep. No rip-and-replace.",
)

# ============ Slide 2: Key Features & Innovation (base: THE PRODUCT) =======
s2 = pres.slides[3]
chip2 = by_text(s2, "THE PRODUCT")
set_text(chip2, "KEY FEATURES & INNOVATION")
chip2.width = Inches(3.6)
for sh in s2.shapes:
    if sh.name == "Shape 0":
        sh.width = Inches(2.62)
set_text(
    by_text(s2, "Root-cause analysis"),
    "Root-cause analysis and guarded medicine transfers between facilities — approved in one click, full audit trail.",
)
set_text(by_text(s2, "Realtime, No Refresh"), "Voice-First Frontline")
set_text(
    by_text(s2, "Every field update"),
    "A worker speaks in Hindi, Marathi or English — Gemini parses it and the whole district re-scores in seconds.",
)
set_text(by_text(s2, "Smart Redistribution"), "WhatsApp Alerts + PDF Reports")
set_text(
    by_text(s2, "Recommends safe transfers"),
    "One-click WhatsApp alerts to the frontline and a full district PDF report — every alert audited live.",
)
# Small "also inside" line under the feature cards.
import copy as _copy

tmpl = by_text(s2, "15 facilities scored")
extra = _copy.deepcopy(tmpl._element)
tmpl._element.getparent().append(extra)
extra_sh = s2.shapes[-1]
extra_sh.left, extra_sh.top = Inches(8.09), Inches(6.86)
extra_sh.width, extra_sh.height = Inches(4.71), Inches(0.5)
set_text(extra_sh, "Also inside: disaster stress mode · deterministic risk engines · 88 unit tests · Gemini + Firestore + Maps.")
for p in extra_sh.text_frame.paragraphs:
    for r in p.runs:
        r.font.color.rgb = MUTED
        r.font.size = Pt(10)

# ========= Slide 3: Impact & Scalability (base: MEASURED IMPACT) ===========
s3 = pres.slides[9]
chip3 = by_text(s3, "MEASURED IMPACT")
set_text(chip3, "IMPACT & SCALABILITY")
chip3.width = Inches(2.4)
for sh in s3.shapes:
    if sh.name == "Shape 1":
        sh.width = Inches(1.62)
set_text(
    by_text(s3, "Counterfactual simulation"),
    "Counterfactual simulation over a 90-day history and 30-day projection — computed by the exact engines and "
    "guardrails that run in the product.",
)

# Scale-and-ask panel over the illustration on the right.
panel = s3.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(8.72), Inches(4.05), Inches(4.25), Inches(3.05))
panel.adjustments[0] = 0.045
panel.fill.solid()
panel.fill.fore_color.rgb = PANEL_BG
panel.line.color.rgb = PANEL_LINE
panel.line.width = Pt(1)
panel.shadow.inherit = False

tb = s3.shapes.add_textbox(Inches(8.98), Inches(4.28), Inches(3.76), Inches(2.62))
tf = tb.text_frame
tf.word_wrap = True

def para(text, *, size, color, bold=False, space_before=6, first=False, hang=False):
    p = tf.paragraphs[0] if first else tf.add_paragraph()
    p.space_before = Pt(space_before)
    if hang:  # wrapped lines align under the text, not under the bullet
        pPr = p._p.get_or_add_pPr()
        pPr.set("marL", "171450")
        pPr.set("indent", "-171450")
    r = p.add_run()
    r.text = text
    r.font.name = "Inter"
    r.font.size = Pt(size)
    r.font.bold = bold
    r.font.color.rgb = color
    return p

para("SCALE ON EXISTING RAILS", size=10.5, color=PINK, bold=True, space_before=0, first=True)
para("•  One district → Maharashtra → 800+ districts, on HMIS & DVDMS rails", size=10.5, color=BODY, hang=True)
para("•  Funded through NHM's Programme Implementation Plan — no new hardware", size=10.5, color=BODY, hang=True)
para("•  Next: offline-first frontline · every regional language · richer readiness scores", size=10.5, color=BODY, hang=True)
para("Our ask: sanction one pilot district. Its own numbers will make the case.", size=11.5, color=PINK, bold=True, space_before=10)

# ============ Keep only the three slides, in order =========================
keep = {2, 3, 9}
sldIdLst = pres.slides._sldIdLst
for idx, sld in reversed(list(enumerate(list(sldIdLst)))):
    if idx not in keep:
        rId = sld.get("{http://schemas.openxmlformats.org/officeDocument/2006/relationships}id")
        pres.part.drop_rel(rId)
        sldIdLst.remove(sld)

pres.save(OUT)
print("WROTE", OUT, "slides:", len(sldIdLst))
